import os
import re
import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


def clean_cell(value):
    if value is None:
        return ""
    return str(value).strip()


def parse_number(value):
    if value is None:
        return None

    text = str(value)
    text = text.replace("$", "").replace(",", "").replace("%", "")
    text = text.replace("M", "").replace("B", "").strip()

    try:
        return float(text)
    except Exception:
        return None


def extract_numbers_from_text(text):
    pattern = r"[-+]?\$?\d[\d,]*(?:\.\d+)?%?"
    matches = re.findall(pattern, text)

    numbers = []

    for match in matches:
        value = parse_number(match)
        if value is not None:
            numbers.append({
                "raw": match,
                "value": value
            })

    return numbers


def dataframe_to_table(df, table_name="Table"):
    df = df.fillna("")
    headers = [str(col).strip() for col in df.columns]

    rows = []
    for _, row in df.iterrows():
        rows.append([clean_cell(v) for v in row.tolist()])

    return {
        "table_name": table_name,
        "headers": headers,
        "rows": rows
    }


def extract_docx(path):
    doc = Document(path)

    text_parts = []
    tables = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())

    for index, table in enumerate(doc.tables):
        rows = []

        for row in table.rows:
            rows.append([clean_cell(cell.text) for cell in row.cells])

        if len(rows) > 1:
            headers = rows[0]
            body = rows[1:]

            tables.append({
                "table_name": f"DOCX Table {index + 1}",
                "headers": headers,
                "rows": body
            })

    text = "\n".join(text_parts)

    return text, tables


def extract_pdf(path):
    reader = PdfReader(path)

    text_parts = []
    tables = []

    for page_index, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        if page_text.strip():
            text_parts.append(f"--- Page {page_index + 1} ---\n{page_text}")

    text = "\n\n".join(text_parts)

    # Basic PDF table fallback:
    # pypdf does not reliably extract tables, so we preserve line-based text.
    # Better table extraction can later use camelot/tabula/pdfplumber.
    return text, tables


def extract_pptx(path):
    prs = Presentation(path)

    text_parts = []
    tables = []

    for slide_index, slide in enumerate(prs.slides):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            if shape.has_table:
                table = shape.table
                rows = []

                for row in table.rows:
                    rows.append([clean_cell(cell.text) for cell in row.cells])

                if len(rows) > 1:
                    tables.append({
                        "table_name": f"PPTX Slide {slide_index + 1} Table",
                        "headers": rows[0],
                        "rows": rows[1:]
                    })

        if slide_text:
            text_parts.append(
                f"--- Slide {slide_index + 1} ---\n" + "\n".join(slide_text)
            )

    return "\n\n".join(text_parts), tables


def extract_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    return text, []


def extract_csv(path):
    df = pd.read_csv(path, skipinitialspace=True, on_bad_lines="skip")
    table = dataframe_to_table(df, "CSV Table")

    return df.to_string(index=False), [table]


def extract_excel(path):
    sheets = pd.read_excel(path, sheet_name=None)

    text_parts = []
    tables = []

    for sheet_name, df in sheets.items():
        table = dataframe_to_table(df, f"Excel Sheet: {sheet_name}")
        tables.append(table)
        text_parts.append(f"--- Sheet: {sheet_name} ---\n{df.to_string(index=False)}")

    return "\n\n".join(text_parts), tables


def build_chart_candidates_from_tables(tables):
    """
    Creates deterministic chart-ready data from extracted tables.
    This reduces hallucination because charts are created from real table values.
    """
    candidates = []

    for table in tables:
        headers = table.get("headers", [])
        rows = table.get("rows", [])

        if not headers or not rows:
            continue

        # Choose first column as label column
        label_col_index = 0

        for col_index in range(1, len(headers)):
            chart_data = []

            for row in rows:
                if len(row) <= col_index:
                    continue

                label = clean_cell(row[label_col_index])
                value = parse_number(row[col_index])

                if label and value is not None:
                    chart_data.append({
                        "label": label,
                        "value": value
                    })

            if len(chart_data) >= 2:
                chart_type = "line" if looks_like_time_series(chart_data) else "bar"

                candidates.append({
                    "title": f"{headers[col_index]} by {headers[label_col_index]}",
                    "chart_type": chart_type,
                    "source_table": table.get("table_name", ""),
                    "data": chart_data[:12]
                })

    return candidates[:8]


def looks_like_time_series(chart_data):
    time_labels = {
        "jan", "feb", "mar", "apr", "may", "jun",
        "jul", "aug", "sep", "oct", "nov", "dec",
        "q1", "q2", "q3", "q4"
    }

    labels = [str(item["label"]).lower() for item in chart_data]

    hits = 0
    for label in labels:
        if label in time_labels:
            hits += 1

    return hits >= 2


def extract_document_data(path):
    """
    Main function.

    Returns:
    {
      text,
      tables,
      numbers,
      chart_candidates
    }
    """
    ext = os.path.splitext(path)[1].lower()

    if ext == ".docx":
        text, tables = extract_docx(path)

    elif ext == ".pdf":
        text, tables = extract_pdf(path)

    elif ext == ".pptx":
        text, tables = extract_pptx(path)

    elif ext == ".txt":
        text, tables = extract_txt(path)

    elif ext == ".csv":
        text, tables = extract_csv(path)

    elif ext in [".xlsx", ".xls"]:
        text, tables = extract_excel(path)

    else:
        raise ValueError(f"Unsupported file type: {ext}")

    numbers = extract_numbers_from_text(text)
    chart_candidates = build_chart_candidates_from_tables(tables)

    return {
        "text": text,
        "tables": tables,
        "numbers": numbers,
        "chart_candidates": chart_candidates
    }