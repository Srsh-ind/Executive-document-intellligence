from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from themes import DARK_BLUE, ACCENT_ORANGE, LIGHT_GRAY, MID_GRAY, BLACK, WHITE


def add_text(slide, text, left, top, width, height, size=14, bold=False, color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = str(text)

    for p in tf.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color

    return box


def add_header(slide, title, page):
    add_text(
        slide,
        title,
        Inches(0.5),
        Inches(0.25),
        Inches(12),
        Inches(0.45),
        22,
        True,
        DARK_BLUE
    )

    line = slide.shapes.add_shape(
        1,
        Inches(0.5),
        Inches(0.82),
        Inches(12.25),
        Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.color.rgb = ACCENT_ORANGE

    add_text(
        slide,
        str(page),
        Inches(12.2),
        Inches(7.05),
        Inches(0.5),
        Inches(0.2),
        9,
        False,
        MID_GRAY
    )


def add_card(slide, title, body, left, top, width, height):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = RGBColor(220, 220, 220)

    add_text(
        slide,
        title,
        left + Inches(0.15),
        top + Inches(0.1),
        width - Inches(0.3),
        Inches(0.35),
        12,
        True,
        DARK_BLUE
    )

    add_text(
        slide,
        body,
        left + Inches(0.15),
        top + Inches(0.5),
        width - Inches(0.3),
        height - Inches(0.55),
        12,
        False,
        BLACK
    )


def add_bullets(slide, items, left, top, width, height, size=15, max_items=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for item in items[:max_items]:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)