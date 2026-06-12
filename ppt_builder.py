import os

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from themes import DARK_BLUE, WHITE, BLACK
from layouts import add_text, add_header, add_card, add_bullets
from slide_planner import build_slide_plan
from chart_agent import create_visual_assets


def clean_bullet_item(item):
    if isinstance(item, str):
        return item

    if isinstance(item, dict):
        if "insight" in item:
            return item.get("insight", "")
        if "description" in item:
            return item.get("description", "")
        if "risk" in item:
            return f"{item.get('risk', '')}: {item.get('description', '')}"
        if "title" in item:
            return item.get("title", "")

    return str(item)


def clean_bullet_list(items):
    cleaned = []

    for item in items:
        text = clean_bullet_item(item)

        if text and not text.strip().startswith("{"):
            cleaned.append(text)

    return cleaned


def create_title_slide(prs, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BLUE

    title = analysis.get("title", "Executive Insight Report")

    add_text(
        slide,
        title,
        Inches(0.75),
        Inches(2.5),
        Inches(12),
        Inches(1.0),
        36,
        True,
        WHITE
    )


def create_summary_slide(prs, analysis, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Executive Summary", page)

    add_text(
        slide,
        analysis.get("executive_summary", ""),
        Inches(0.8),
        Inches(1.25),
        Inches(12),
        Inches(4.9),
        18,
        False,
        BLACK
    )


def create_metrics_slide(prs, analysis, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Key Metrics & Signals", page)

    metrics = analysis.get("metrics", [])

    if not metrics:
        metrics = analysis.get("key_metrics", [])

    if not metrics:
        add_text(
            slide,
            "No explicit numeric metrics were detected. The analysis focuses on qualitative themes and implications.",
            Inches(0.8),
            Inches(1.4),
            Inches(11.5),
            Inches(1),
            16
        )
        return

    positions = [
        (0.7, 1.3),
        (4.7, 1.3),
        (8.7, 1.3),
        (0.7, 3.2),
        (4.7, 3.2),
        (8.7, 3.2)
    ]

    for idx, metric in enumerate(metrics[:6]):
        x, y = positions[idx]

        if isinstance(metric, dict):
            title = f"{metric.get('name', 'Metric')}: {metric.get('value', '')}"
            body = metric.get("interpretation", "")
        else:
            title = "Metric"
            body = str(metric)

        add_card(
            slide,
            title,
            body,
            Inches(x),
            Inches(y),
            Inches(3.6),
            Inches(1.45)
        )


def create_chart_slide(prs, chart_item, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, chart_item.get("title", "Business Visualization"), page)

    slide.shapes.add_picture(
        chart_item["path"],
        Inches(0.65),
        Inches(1.15),
        width=Inches(7.6)
    )

    add_card(
        slide,
        "Executive Interpretation",
        chart_item.get(
            "insight",
            "This visualization highlights a key pattern extracted from the source document."
        ),
        Inches(8.55),
        Inches(1.35),
        Inches(4.05),
        Inches(3.6)
    )


def create_bullet_slide(prs, title, items, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, page)

    cleaned_items = clean_bullet_list(items)

    if not cleaned_items:
        cleaned_items = ["No material items were identified for this section."]

    add_bullets(
        slide,
        cleaned_items,
        Inches(0.8),
        Inches(1.25),
        Inches(11.8),
        Inches(5.6),
        15,
        7
    )


def create_recommendation_slide(prs, analysis, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Strategic Recommendations", page)

    recs = analysis.get("recommendations", [])

    if not recs:
        add_text(
            slide,
            "No explicit recommendations were generated for this document.",
            Inches(0.8),
            Inches(1.4),
            Inches(11.5),
            Inches(1),
            16
        )
        return

    top = Inches(1.4)

    for rec in recs[:4]:
        if isinstance(rec, dict):
            priority = rec.get("priority", "Medium")
            recommendation = rec.get("recommendation", "")
            impact = rec.get("business_impact", "")
        else:
            priority = "Medium"
            recommendation = str(rec)
            impact = ""

        body = f"{recommendation}\n\nImpact: {impact}"

        add_card(
            slide,
            f"{priority} Priority",
            body,
            Inches(0.85),
            top,
            Inches(11.6),
            Inches(1.35)
        )

        top += Inches(1.55)


def create_conclusion_slide(prs, analysis, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Executive Conclusion", page)

    conclusion = analysis.get("conclusion", "")

    add_text(
        slide,
        conclusion,
        Inches(0.8),
        Inches(1.35),
        Inches(11.8),
        Inches(4.8),
        19,
        False,
        BLACK
    )


def create_ppt(analysis):
    os.makedirs("outputs", exist_ok=True)

    analysis = create_visual_assets(analysis)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, analysis)

    page = 2

    slide_plan = build_slide_plan(analysis)

    for slide_def in slide_plan:
        slide_type = slide_def.get("type")
        title = slide_def.get("title")

        if slide_type == "summary":
            create_summary_slide(prs, analysis, page)
            page += 1

        elif slide_type == "metrics":
            create_metrics_slide(prs, analysis, page)
            page += 1

            for chart_item in analysis.get("chart_paths", []):
                create_chart_slide(prs, chart_item, page)
                page += 1

        elif slide_type == "charts":
            if "metrics" not in [s.get("type") for s in slide_plan]:
                for chart_item in analysis.get("chart_paths", []):
                    create_chart_slide(prs, chart_item, page)
                    page += 1

        elif slide_type == "bullets":
            items = slide_def.get("items")

            if items is None:
                source = slide_def.get("source", "")
                items = analysis.get(source, [])

            create_bullet_slide(prs, title, items, page)
            page += 1

        elif slide_type == "custom_section":
            create_bullet_slide(
                prs,
                title,
                slide_def.get("items", []),
                page
            )
            page += 1

        elif slide_type == "recommendations":
            create_recommendation_slide(prs, analysis, page)
            page += 1

        elif slide_type == "conclusion":
            create_conclusion_slide(prs, analysis, page)
            page += 1

    ppt_file = "outputs/executive_document_intelligence.pptx"
    prs.save(ppt_file)

    return ppt_file