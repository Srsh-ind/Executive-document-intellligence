from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
import re

from presentation.layout_utils import (
    safe_box, add_picture_safe, SAFE_BOTTOM, SLIDE_W, SLIDE_H, FOOTER_TOP
)
from presentation.widgets import (
    add_background,
    add_title,
    add_subtitle,
    add_textbox,
    add_multiline_textbox,
    add_card,
    add_kpi_card,
    add_metric_strip,
    add_status_card,
    add_bullet_list,
    add_footer,
    add_progress_bar,
    add_section_label,
    add_header_band,
    add_divider_line,
    rgb,
    fit_text,
)


# ─────────────────────────── helpers ────────────────────────────

def get_chart_for_slide(analysis, index=0):
    charts = analysis.get("chart_paths", []) or []
    return charts[index] if len(charts) > index else None


def get_charts_of_type(analysis, chart_types):
    charts = analysis.get("chart_paths", []) or []
    return [c for c in charts if str(c.get("chart_type", "")).lower() in chart_types]


def add_chart_image(slide, chart, x, y, w, h=None, reserved_bottom=None):
    """Add chart image, auto-clamping height so it never overflows the slide."""
    if not chart:
        return False
    path = chart.get("path")
    if not path:
        return False
    return add_picture_safe(slide, path, x, y, w, h=h, reserved_bottom=reserved_bottom)


def get_cards(slide):    return slide.get("cards", []) or []
def get_bullets(slide):  return slide.get("bullets", []) or []
def get_blocks(slide):   return slide.get("blocks", []) or []


def _display_clean(text):
    """Remove evidence/debug labels from content shown on slides."""
    text = str(text or "").strip()
    lines = []
    for line in text.splitlines():
        low = line.strip().lower()
        if low.startswith(("audience:", "domain expected:", "prepared for:", "reporting period:", "source profile:")):
            continue
        lines.append(line)
    text = "\n".join(lines).strip()
    text = re.sub(r"\s*\|?\s*Evidence\s*(?:IDs?|references?)?\s*[:：]\s*(?:E\d{3}(?:\s*[,;]\s*)?)+", "", text, flags=re.I)
    text = re.sub(r"\s*\(?\bE\d{3}(?:\s*[,;]\s*E\d{3})*\)?", "", text, flags=re.I)
    text = re.sub(r"^(Implication|Reasoning|Business implication)\s*[:：]\s*", "", text, flags=re.I)
    text = text.replace("Evidence → interpretation → implication:", "")
    text = text.replace("Evidence -> interpretation -> implication:", "")
    text = text.replace("Evidence → interpretation → action → expected impact:", "")
    text = text.replace("Evidence -> interpretation -> action -> expected impact:", "")
    text = re.sub(r"\bStrategies should focus on\s*$", "", text, flags=re.I)
    text = text.replace("…", "")
    text = re.sub(r"\.{3,}", "", text)
    return re.sub(r"\s+", " ", text).strip()


def _visible_reasoning(text):
    text = _display_clean(text)
    low = text.lower()
    if not text or any(bad in low for bad in (
        "retrieved supporting evidence",
        "cited evidence",
        "evidence item",
        "directly derived",
        "contains a numeric document fact",
        "reasoning path is",
    )):
        return ""
    return text[:360].rstrip()


def card_text(card):
    if isinstance(card, dict):
        label = (card.get("label") or card.get("name") or card.get("priority")
                 or card.get("title") or card.get("risk") or card.get("opportunity") or "Insight")
        value = (card.get("value") or card.get("recommendation") or card.get("action")
                 or card.get("risk") or card.get("opportunity") or card.get("item") or card.get("claim") or "")
        note  = (card.get("note") or card.get("reasoning") or card.get("interpretation") or card.get("description")
                 or card.get("business_implication") or card.get("business_impact") or card.get("impact") or "")
        return _display_clean(label), _display_clean(value), _visible_reasoning(note)
    return "Insight", _display_clean(card), ""


def trim_list(values, limit):
    return list(values)[:limit] if values else []


def add_decorative_circle(slide, theme, x=10.6, y=0.1, size=3.2, alpha_color=None):
    colors = theme["colors"]
    color = alpha_color or colors["accent"]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_decorative_rect(slide, theme, x, y, w, h, color_key="primary", alpha=None):
    colors = theme["colors"]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(alpha or colors[color_key])
    shape.line.fill.background()
    return shape


# ────────────────────────── cover ───────────────────────────────

def render_hero_cover(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    # Bold accent circle top-right — fully contained within slide
    _circle_size = 4.8
    _big_x = SLIDE_W - _circle_size - 0.1   # flush right, no overflow
    _big_y = 0.0                              # top of slide, no bleed
    big = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(_big_x), Inches(_big_y), Inches(_circle_size), Inches(_circle_size)
    )
    big.fill.solid()
    big.fill.fore_color.rgb = rgb(colors["primary"])
    big.line.fill.background()

    # Smaller accent circle bottom-right — clamped to slide bounds
    _sm_size = 2.4
    _sm_x = SLIDE_W - _sm_size - 0.1
    _sm_y = min(4.8, SLIDE_H - _sm_size - 0.1)
    small = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(_sm_x), Inches(_sm_y), Inches(_sm_size), Inches(_sm_size)
    )
    small.fill.solid()
    small.fill.fore_color.rgb = rgb(colors["accent"])
    small.line.fill.background()

    # Accent bar at bottom
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.22), Inches(13.33), Inches(0.28)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(colors["accent"])
    bar.line.fill.background()

    doc_type = analysis.get("document_type", "Executive Report").upper()
    add_textbox(slide, doc_type, 0.7, 1.5, 8.5, 0.4,
                font_size=13, bold=True, color=colors["accent"])

    title_text = slide_data.get("title") or analysis.get("title", "Executive Insight Report")
    # Limit title width to 7.8" so it never overlaps the circle (circle starts at ~8.43")
    add_textbox(slide, title_text, 0.7, 2.0, 7.6, 2.2,
                font_size=36, bold=True, color=colors["text"])

    subtitle = slide_data.get("headline") or analysis.get("core_message", "")
    add_textbox(slide, subtitle, 0.7, 4.3, 7.4, 1.0,
                font_size=16, bold=False, color=colors["subtext"])

    # Do not render audience/persona tags on the cover; they looked like stray UI pills
    # (e.g. "CEO, CFO, COO...") in generated decks.
    add_footer(slide, theme, page)


# ────────────────────────── summary_dashboard ───────────────────

def render_summary_dashboard(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    add_header_band(slide, theme)
    colors = theme["colors"]

    # No decorative circle on summary — the metric strip at y=1.52 conflicts; skip it
    add_title(slide, slide_data.get("title", "Executive Summary"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)

    # Metric strip from top metrics.  Give cards more width/height than before
    # so notes such as YoY growth / margin drivers stay inside the upper boxes.
    metric_cards = get_cards(slide_data) or trim_list(analysis.get("metrics"), 4)
    if metric_cards:
        add_metric_strip(slide, metric_cards[:4], theme, x=0.45, y=1.44, w=12.43, h=1.40)

    blocks = get_blocks(slide_data)
    if not blocks:
        blocks = [
            {"label": "Key Findings",   "items": trim_list(analysis.get("key_findings"), 4)},
            {"label": "Insights",       "items": trim_list(analysis.get("insights"), 4)},
            {"label": "Next Steps",     "items": [card_text(r)[1] or card_text(r)[0]
                                                  for r in trim_list(analysis.get("recommendations"), 3)]}
        ]
        if not blocks[1]["items"] and analysis.get("opportunities"):
            blocks[1] = {"label": "Opportunities", "items": trim_list(analysis.get("opportunities"), 4)}

    # Lower boxes are intentionally lower and shorter than the metric strip;
    # this prevents long top-card notes from visually colliding with the blocks.
    positions = [(0.48, 3.24), (4.68, 3.24), (8.88, 3.24)]
    for i, block in enumerate(blocks[:3]):
        x, y = positions[i]
        items_text = "\n".join([_display_clean(str(it)) for it in (block.get("items") or [])[:3] if _display_clean(str(it))])
        add_card(slide, block.get("label", "Section"), items_text, theme, x, y, 3.92, 3.10)

    add_footer(slide, theme, page)


# ────────────────────────── kpi_dashboard ───────────────────────

def render_kpi_dashboard(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    add_header_band(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "KPI Dashboard"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)

    cards = [c for c in (get_cards(slide_data) or trim_list(analysis.get("metrics"), 6)) if card_text(c)[1]]
    if len(cards) < 2:
        # Do not render a mostly empty KPI slide; summary and insight slides carry the metric.
        add_title(slide, slide_data.get("title", "Performance Snapshot"), theme)
        add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
        bullets = trim_list(analysis.get("insights") or analysis.get("key_findings"), 4)
        add_bullet_list(slide, bullets, theme, 0.7, 1.65, 11.8, 5.2, font_size=13)
        add_footer(slide, theme, page)
        return
    # Accent bar decoration
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    positions_2rows = [
        (0.55, 1.60), (4.55, 1.60), (8.55, 1.60),
        (0.55, 3.70), (4.55, 3.70), (8.55, 3.70),
    ]
    # 4-card row: total = 4*kw + 3*gap + 2*margin
    # kw=2.95, gap=0.45, margin=0.55 → 4*2.95 + 3*0.45 + 2*0.55 = 11.8+1.35+1.1 = 13.25 ✓
    positions_1row = [
        (0.55, 2.30), (3.55, 2.30), (6.55, 2.30), (9.55, 2.30),
    ]

    count = min(len(cards), 6)
    if count <= 4:
        pos = positions_1row
        kw, kh = 2.95, 2.6
    else:
        pos = positions_2rows
        kw, kh = 3.7, 1.75

    for i, card in enumerate(cards[:count]):
        label, value, note = card_text(card)
        x, y = pos[i]
        add_kpi_card(slide, label, value, note, theme, x, y, kw, kh)

    add_footer(slide, theme, page)


# ────────────────────────── performance_dashboard ───────────────
# Shows charts — distribute ALL chart_paths across chart slides

def render_performance_dashboard(prs, slide_data, analysis, theme, page):
    """
    Chart-focused slide. Shows up to 4 charts in a 2×2 grid,
    or 2 side-by-side, or 1 large.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    charts = analysis.get("chart_paths", []) or []
    # Use a running index stored on analysis to distribute charts across slides
    _chart_cursor = analysis.get("_perf_chart_cursor", 0)
    charts_per_slide = 4
    window = charts[_chart_cursor:_chart_cursor + charts_per_slide]
    if not window:
        window = charts[:charts_per_slide]
        analysis["_perf_chart_cursor"] = charts_per_slide
    else:
        analysis["_perf_chart_cursor"] = _chart_cursor + charts_per_slide

    # Auto-generate a meaningful title from the chart titles in this window
    _default_title = slide_data.get("title", "Performance Trends")
    if window and _default_title in ("Performance Trends", "Data Insights II", "Data Insights III",
                                     "Data Insights IV", "Data Insights V", "Additional Data Insights",
                                     "Further quantitative evidence"):
        _chart_titles = [c.get("title", "") if isinstance(c, dict) else "" for c in window]
        _chart_titles = [t for t in _chart_titles if t]
        if _chart_titles:
            if len(_chart_titles) == 1:
                _auto_title = _chart_titles[0]
            elif len(_chart_titles) == 2:
                _auto_title = f"{_chart_titles[0]} & {_chart_titles[1]}"
            else:
                _auto_title = f"{_chart_titles[0]}, {_chart_titles[1]} & More"
            _slide_title = _auto_title
        else:
            _slide_title = _default_title
    else:
        _slide_title = _default_title

    add_title(slide, _slide_title, theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    n = len(window)
    if n >= 4:
        # 2×2 grid
        positions = [
            (0.55, 1.58, 6.1), (6.78, 1.58, 6.1),
            (0.55, 4.10, 6.1), (6.78, 4.10, 6.1),
        ]
        for i, chart in enumerate(window[:4]):
            x, y, w = positions[i]
            # max height for each cell so we stay inside slide
            _rb = FOOTER_TOP - 0.08 if y >= 3.5 else (positions[2][1] - 0.08)
            add_chart_image(slide, chart, x, y, w, h=2.20, reserved_bottom=_rb)
    elif n == 3:
        # Use a 2x2 grid with three equal-sized chart cells; do not enlarge the third chart.
        positions = [
            (0.55, 1.58, 6.1), (6.78, 1.58, 6.1),
            (0.55, 4.10, 6.1),
        ]
        for i, chart in enumerate(window[:3]):
            x, y, w = positions[i]
            _rb = (positions[2][1] - 0.08) if y < 3.5 else (FOOTER_TOP - 0.08)
            add_chart_image(slide, chart, x, y, w, h=2.20, reserved_bottom=_rb)
    elif n == 2:
        add_chart_image(slide, window[0], 0.55, 1.62, 6.1, h=3.25)
        add_chart_image(slide, window[1], 6.78, 1.62, 6.1, h=3.25)
        so_what = slide_data.get("so_what") or slide_data.get("message", "")
        if so_what:
            add_card(slide, "Key Takeaway", so_what, theme, 0.55, 5.22, 12.2, 1.55)
    elif n == 1:
        add_chart_image(slide, window[0], 0.55, 1.62, 7.5, h=4.25)
        so_what = slide_data.get("so_what") or slide_data.get("message", "")
        add_card(slide, "Interpretation", so_what, theme, 8.25, 1.62, 4.4, 5.4)
    else:
        # No charts — fall back to insight cards
        bullets = get_bullets(slide_data) or trim_list(analysis.get("key_findings"), 5)
        add_bullet_list(slide, bullets, theme, 0.7, 1.65, 11.8, 5.2, font_size=14)

    add_footer(slide, theme, page)


# ────────────────────────── split_metrics_chart ─────────────────

def render_split_metrics_chart(prs, slide_data, analysis, theme, page):
    """Left: KPI cards. Right: chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Metrics & Trends"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    # Left: up to 3 KPI cards stacked
    cards = get_cards(slide_data) or trim_list(analysis.get("metrics"), 3)
    y = 1.62
    for card in cards[:3]:
        label, value, note = card_text(card)
        add_kpi_card(slide, label, value, note, theme, 0.55, y, 5.5, 1.65)
        y += 1.82

    # Right: chart
    charts = analysis.get("chart_paths", []) or []
    chart = charts[0] if charts else None
    if chart:
        add_chart_image(slide, chart, 6.35, 1.58, 6.6)
    else:
        add_card(slide, "Interpretation", slide_data.get("so_what", ""), theme,
                 6.35, 1.58, 6.6, 5.4)

    add_footer(slide, theme, page)


# ────────────────────────── insight_dashboard ───────────────────

def render_insight_dashboard(prs, slide_data, analysis, theme, page):
    """Insight cards — 3 or 6 cards with coloured icon-number accent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Key Insights"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    bullets = get_bullets(slide_data)
    cards   = get_cards(slide_data)

    # Fallback to analysis data
    if not cards and not bullets:
        role = slide_data.get("story_role", "")
        if "risk" in role:
            cards = trim_list(analysis.get("risks"), 6)
        elif "opportun" in role:
            cards = trim_list(analysis.get("opportunities"), 6)
        elif "action" in role or "recommend" in role:
            cards = trim_list(analysis.get("recommendations"), 6)
        else:
            bullets = trim_list(analysis.get("insights") or analysis.get("key_findings"), 6)

    if cards:
        count = min(len(cards), 6)
        if count <= 3:
            pos = [(0.55, 1.65), (4.55, 1.65), (8.55, 1.65)]
            cw, ch = 3.75, 5.15
        else:
            pos = [
                (0.55, 1.65), (4.55, 1.65), (8.55, 1.65),
                (0.55, 3.85), (4.55, 3.85), (8.55, 3.85),
            ]
            cw, ch = 3.75, 2.0

        for i, card in enumerate(cards[:count]):
            label, value, note = card_text(card)
            body = f"{value}\n{note}".strip() if note else value
            x, y = pos[i]
            # Number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 0.06), Inches(y + 0.1),
                Inches(0.42), Inches(0.42)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = rgb(colors["accent"])
            badge.line.fill.background()
            add_textbox(slide, str(i + 1), x + 0.06, y + 0.12,
                        0.42, 0.28, font_size=11, bold=True, color="FFFFFF")

            add_card(slide, label, body, theme, x, y, cw, ch)
    elif bullets:
        # Numbered list with accent circles
        _lst_end  = SAFE_BOTTOM - 0.08
        _lst_avail = _lst_end - 1.65
        _n_blt    = min(len(bullets), 6)
        step = min(_lst_avail / max(_n_blt, 1), 0.90)
        y = 1.65
        for i, item in enumerate(bullets[:_n_blt]):
            if y + step > _lst_end:
                break
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.55), Inches(y + 0.05),
                Inches(0.38), Inches(0.38)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = rgb(colors["accent"])
            badge.line.fill.background()
            add_textbox(slide, str(i + 1), 0.55, y + 0.08,
                        0.38, 0.24, font_size=9, bold=True, color="FFFFFF")
            add_textbox(slide, str(item), 1.12, y, 11.2, step - 0.04,
                        font_size=13, bold=False, color=colors["text"])
            _blt_div = min(y + step - 0.04, _lst_end)
            add_divider_line(slide, theme, 1.1, _blt_div, 11.6)
            y += step

    add_footer(slide, theme, page)


# ────────────────────────── three_column_status ─────────────────

def render_three_column_status(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Status Overview"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    blocks = get_blocks(slide_data)
    if not blocks:
        blocks = [
            {"label": "Strengths",  "items": trim_list(analysis.get("key_findings"), 5)},
            {"label": "Concerns",   "items": trim_list(analysis.get("risks"), 5)},
            {"label": "Priorities", "items": [card_text(r)[1] or card_text(r)[0]
                                              for r in trim_list(analysis.get("recommendations"), 4)]},
        ]

    col_colors = [colors["success"], colors["danger"], colors["accent"]]
    positions = [(0.55, 1.60), (4.72, 1.60), (8.89, 1.60)]

    for i, block in enumerate(blocks[:3]):
        x, y = positions[i]
        # Color accent bar at top of each column
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.72), Inches(0.10)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = rgb(col_colors[i % len(col_colors)])
        accent_bar.line.fill.background()

        items = []
        for item in (block.get("items") or [])[:5]:
            if isinstance(item, dict):
                t = card_text(item)
                items.append(t[1] or t[0])
            else:
                items.append(str(item))

        # Height clamped: y+0.10 + h must not exceed SAFE_BOTTOM (7.06)
        _card_y = y + 0.10
        _card_h = min(5.6, SAFE_BOTTOM - _card_y - 0.04)
        add_card(slide, block.get("label", "Section"),
                 "\n".join(items), theme, x, _card_y, 3.72, _card_h)

    add_footer(slide, theme, page)


# ────────────────────────── risk_dashboard ──────────────────────

def render_risk_dashboard(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Risks & Challenges"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    cards = get_cards(slide_data)
    if not cards:
        cards = trim_list(analysis.get("risks"), 3) + trim_list(analysis.get("opportunities"), 3)

    positions = [
        (0.55, 1.62), (4.55, 1.62), (8.55, 1.62),
        (0.55, 3.95), (4.55, 3.95), (8.55, 3.95),
    ]
    role = slide_data.get("story_role", "")

    if cards:
        for i, card in enumerate(cards[:6]):
            label, value, note = card_text(card)
            status = "risk" if "risk" in role else "success" if "opportun" in role else "neutral"
            add_status_card(slide, label, value, note, theme,
                            positions[i][0], positions[i][1], 3.75, 2.0, status=status)
    else:
        bullets = get_bullets(slide_data) or trim_list(analysis.get("key_findings") or analysis.get("insights"), 6)
        add_bullet_list(slide, bullets, theme, 0.7, 1.72, 11.8, 5.1, font_size=13)

    add_footer(slide, theme, page)


# ────────────────────────── opportunity_dashboard ───────────────

def render_opportunity_dashboard(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Opportunities"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    cards = get_cards(slide_data) or trim_list(analysis.get("opportunities"), 6)

    positions = [
        (0.55, 1.62), (4.55, 1.62), (8.55, 1.62),
        (0.55, 3.95), (4.55, 3.95), (8.55, 3.95),
    ]
    if cards:
        for i, card in enumerate(cards[:6]):
            label, value, note = card_text(card)
            add_status_card(slide, label, value, note, theme,
                            positions[i][0], positions[i][1], 3.75, 2.0, status="success")
    else:
        bullets = get_bullets(slide_data) or trim_list(analysis.get("opportunities") or analysis.get("key_findings"), 6)
        add_bullet_list(slide, bullets, theme, 0.7, 1.72, 11.8, 5.1, font_size=13)

    add_footer(slide, theme, page)


# ────────────────────────── action_tracker ──────────────────────

def render_action_tracker(prs, slide_data, analysis, theme, page):
    """
    Recommendations as horizontal action rows — priority label + text + impact.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Recommended Actions"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    cards = get_cards(slide_data) or trim_list(analysis.get("recommendations"), 5)
    if not cards and get_bullets(slide_data):
        cards = [{"label": "Action", "value": b, "note": ""} for b in get_bullets(slide_data)[:5]]

    # Column headers
    add_textbox(slide, "PRIORITY", 0.55, 1.56, 1.4, 0.28,
                font_size=9, bold=True, color=colors["subtext"])
    add_textbox(slide, "ACTION", 2.10, 1.56, 6.5, 0.28,
                font_size=9, bold=True, color=colors["subtext"])
    add_textbox(slide, "IMPACT", 8.75, 1.56, 3.9, 0.28,
                font_size=9, bold=True, color=colors["subtext"])
    add_divider_line(slide, theme, 0.55, 1.84, 12.2, accent=True)

    y = 1.94
    row_h = min(0.96, (SAFE_BOTTOM - 0.10 - 1.94) / 5)
    for i, card in enumerate(cards[:5]):
        label, value, note = card_text(card)
        priority = label if label not in ("Insight",) else ("High" if i == 0 else "Medium")

        # Priority badge
        p_color = colors["danger"] if "high" in priority.lower() else \
                  colors["accent"] if "med" in priority.lower() else colors["success"]
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.55), Inches(y + 0.18), Inches(1.4), Inches(0.42)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*[int(p_color[j:j+2], 16) for j in (0, 2, 4)])
        badge.line.fill.background()
        add_textbox(slide, priority, 0.55, y + 0.21, 1.4, 0.30,
                    font_size=10, bold=True, color="FFFFFF")

        add_textbox(slide, fit_text(value, 180), 2.10, y + 0.08, 6.5, row_h - 0.16,
                    font_size=13, bold=False, color=colors["text"])
        add_textbox(slide, fit_text(note, 150), 8.75, y + 0.08, 3.9, row_h - 0.16,
                    font_size=11, bold=False, color=colors["subtext"])
        _div_y = min(y + row_h - 0.04, SAFE_BOTTOM - 0.02)
        add_divider_line(slide, theme, 0.55, _div_y, 12.2)
        y += row_h

    add_footer(slide, theme, page)


# ────────────────────────── roadmap ─────────────────────────────

def render_roadmap(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Roadmap"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)

    # Horizontal timeline line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(3.55), Inches(10.9), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(colors["accent"])
    line.line.fill.background()

    steps = get_cards(slide_data) or get_bullets(slide_data) or trim_list(analysis.get("recommendations"), 4)
    if not steps:
        steps = ["Assess & Plan", "Execute", "Scale", "Measure Impact"]

    count = min(len(steps), 4)
    # Evenly space nodes
    x_positions = [1.2 + i * (10.5 / max(count - 1, 1)) for i in range(count)]
    if count == 1:
        x_positions = [6.2]

    for i, step in enumerate(steps[:count]):
        x = x_positions[i]

        # Circle node
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.38), Inches(3.22), Inches(0.76), Inches(0.76)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(colors["accent"])
        circle.line.color.rgb = rgb(colors["background"])
        circle.line.width = Emu(19050)
        add_textbox(slide, str(i + 1), x - 0.38, 3.30, 0.76, 0.44,
                    font_size=13, bold=True, color="FFFFFF")

        if isinstance(step, dict):
            step_title = step.get("label") or step.get("priority") or f"Step {i+1}"
            step_body  = step.get("value") or step.get("recommendation") or step.get("note", "")
        else:
            step_title = f"Step {i + 1}"
            step_body  = str(step)

        card_x = max(x - 1.45, 0.1)
        card_w = min(2.9, 13.1 - card_x)

        # Alternate cards above/below the line
        if i % 2 == 0:
            add_card(slide, step_title, step_body, theme, card_x, 4.18, card_w, 2.35)
        else:
            add_card(slide, step_title, step_body, theme, card_x, 1.58, card_w, 1.55)

    add_footer(slide, theme, page)


# ────────────────────────── comparison_dashboard ────────────────

def render_comparison_dashboard(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Comparison"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    cards = get_cards(slide_data)

    # Vertical divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(1.6), Inches(0.04), Inches(5.4)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = rgb(colors["accent"])
    div.line.fill.background()

    if len(cards) >= 2:
        l_label, l_value, l_note = card_text(cards[0])
        r_label, r_value, r_note = card_text(cards[1])

        add_kpi_card(slide, l_label, l_value, l_note, theme, 0.55, 1.65, 5.75, 3.5)
        add_kpi_card(slide, r_label, r_value, r_note, theme, 6.78, 1.65, 5.75, 3.5)

        # Additional cards if present
        for i, extra in enumerate(cards[2:4]):
            el, ev, en = card_text(extra)
            x = 0.55 if i % 2 == 0 else 6.78
            add_card(slide, el, f"{ev}\n{en}".strip(), theme, x, 5.35, 5.75, 1.4)
    else:
        bullets = get_bullets(slide_data) or trim_list(analysis.get("key_findings"), 6)
        add_bullet_list(slide, bullets[:3], theme, 0.7, 1.65, 5.55, 5.3, font_size=13)
        add_bullet_list(slide, bullets[3:6], theme, 6.78, 1.65, 5.55, 5.3, font_size=13)

    add_footer(slide, theme, page)


# ────────────────────────── agenda_cards ────────────────────────

def render_agenda_cards(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Agenda"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    items = []
    for block in get_blocks(slide_data):
        items.extend(block.get("items", []))
    if not items:
        items = get_bullets(slide_data) or trim_list(analysis.get("sections"), 6)

    _alist_end = SAFE_BOTTOM - 0.08
    _a_avail   = _alist_end - 1.68
    _a_n       = min(len(items), 6)
    step = min(_a_avail / max(_a_n, 1), 0.88)
    y = 1.68
    for i, item in enumerate(items[:_a_n]):
        if y + step > _alist_end:
            break
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.55), Inches(y + 0.10),
            Inches(0.52), Inches(0.52)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(colors["accent"])
        circle.line.fill.background()
        add_textbox(slide, f"{i+1:02d}", 0.55, y + 0.16, 0.52, 0.28,
                    font_size=12, bold=True, color="FFFFFF")
        add_textbox(slide, str(item), 1.25, y + 0.08, 11.0, step - 0.12,
                    font_size=15, bold=False, color=colors["text"])
        _adiv_y = min(y + step - 0.02, _alist_end)
        add_divider_line(slide, theme, 1.2, _adiv_y, 11.1)
        y += step

    add_footer(slide, theme, page)


# ────────────────────────── donut_insights ──────────────────────
# BUG FIX: was ignoring chart data from chart_paths; now prioritises
# donut/pie charts and falls back gracefully.

def render_donut_insights(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Composition & Share"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    # Prefer donut/pie charts; fall back to any chart
    donut_charts = get_charts_of_type(analysis, {"donut", "doughnut", "pie"})
    other_charts  = [c for c in (analysis.get("chart_paths", []) or [])
                     if c not in donut_charts]
    all_charts = donut_charts + other_charts

    if len(all_charts) >= 2:
        add_chart_image(slide, all_charts[0], 0.55, 1.58, 6.1)
        add_chart_image(slide, all_charts[1], 6.68, 1.58, 6.1)
        so_what = slide_data.get("so_what") or slide_data.get("message", "")
        if so_what:
            add_card(slide, "What This Means", so_what, theme,
                     0.55, 5.52, 12.2, 1.52)
    elif len(all_charts) == 1:
        add_chart_image(slide, all_charts[0], 0.55, 1.62, 6.8)
        # KPI cards on the right
        kpi_cards = get_cards(slide_data) or trim_list(analysis.get("metrics"), 3)
        kpi_y = 1.62
        for card in kpi_cards[:3]:
            label, value, note = card_text(card)
            add_kpi_card(slide, label, value, note, theme, 7.55, kpi_y, 5.2, 1.58)
            kpi_y += 1.74
    else:
        # No charts at all — show metrics as big KPI cards
        cards = get_cards(slide_data) or trim_list(analysis.get("metrics"), 6)
        positions = [
            (0.55, 1.62), (4.55, 1.62), (8.55, 1.62),
            (0.55, 3.72), (4.55, 3.72), (8.55, 3.72),
        ]
        for i, card in enumerate(cards[:6]):
            label, value, note = card_text(card)
            add_kpi_card(slide, label, value, note, theme,
                         positions[i][0], positions[i][1], 3.75, 1.75)

    add_footer(slide, theme, page)


# ────────────────────────── closing_slide ───────────────────────

def render_closing_slide(prs, slide_data, analysis, theme, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    # Accent bottom band
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.22), Inches(13.33), Inches(0.28)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(colors["accent"])
    bar.line.fill.background()

    # Oval clamped so it stays inside slide (y + size <= SLIDE_H)
    # circle fully inside slide: right = 9.9+3.2=13.1 ✓; bottom = y+3.2 <= 7.1 (above footer)
    add_decorative_circle(slide, theme, x=9.9, y=min(3.9, SLIDE_H - 3.2 - 0.4), size=3.2)

    headline   = slide_data.get("headline") or slide_data.get("title", "Thank You")
    so_what    = slide_data.get("so_what")   or slide_data.get("message", "")
    conclusion = analysis.get("conclusion", "")

    # Use the richer of (headline, so_what, conclusion) — avoid duplication
    body_text = conclusion or so_what or headline

    add_textbox(slide, headline, 0.9, 2.1, 9.2, 1.5,
                font_size=30, bold=True, color=colors["text"])
    if body_text and body_text.strip() != headline.strip():
        # Single body paragraph — no duplicate
        add_textbox(slide, body_text, 0.95, 3.80, 8.8, 2.8,
                    font_size=14, bold=False, color=colors["subtext"])

    add_footer(slide, theme, page)


# ────────────────────────── dashboard_grid ──────────────────────

def render_dashboard_grid(prs, slide_data, analysis, theme, page):
    """Mixed grid: 2 KPI cards top + 1 chart bottom-left + insight card bottom-right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme)
    colors = theme["colors"]

    add_title(slide, slide_data.get("title", "Dashboard"), theme)
    add_subtitle(slide, slide_data.get("headline") or slide_data.get("message", ""), theme)
    add_divider_line(slide, theme, 0.55, 1.48, 12.2, accent=True)

    cards = get_cards(slide_data) or trim_list(analysis.get("metrics"), 4)

    # Top row: 2–4 KPI cards
    row_cards = cards[:4]
    count = len(row_cards)
    kw = (12.2 - 0.15 * (count - 1)) / max(count, 1)
    for i, card in enumerate(row_cards):
        label, value, note = card_text(card)
        add_kpi_card(slide, label, value, note, theme,
                     0.55 + i * (kw + 0.15), 1.60, kw, 1.25)

    # Bottom: chart left + (chart or insight) right
    charts = analysis.get("chart_paths", []) or []
    _cursor = analysis.get("_grid_chart_cursor", 0)
    if len(charts) > _cursor + 1:
        # Two charts side-by-side if we have them
        add_chart_image(slide, charts[_cursor], 0.55, 3.25, 6.1, h=3.15)
        add_chart_image(slide, charts[_cursor + 1], 6.78, 3.25, 6.1, h=3.15)
        analysis["_grid_chart_cursor"] = _cursor + 2
    elif len(charts) > _cursor:
        add_chart_image(slide, charts[_cursor], 0.55, 3.25, 7.6, h=3.15)
        analysis["_grid_chart_cursor"] = _cursor + 1
        so_what = slide_data.get("so_what") or slide_data.get("message", "")
        add_card(slide, "Key Takeaway", so_what, theme, 8.35, 3.25, 4.4, 3.15)
    else:
        so_what = slide_data.get("so_what") or slide_data.get("message", "")
        add_card(slide, "Key Takeaway", so_what, theme, 0.55, 3.25, 12.2, 3.15)

    add_footer(slide, theme, page)


# ────────────────────────── before_after ────────────────────────

def render_before_after(prs, slide_data, analysis, theme, page):
    return render_comparison_dashboard(prs, slide_data, analysis, theme, page)


# ──────────────────────── recommendation_roadmap ────────────────

def render_recommendation_roadmap(prs, slide_data, analysis, theme, page):
    return render_roadmap(prs, slide_data, analysis, theme, page)


# ────────────────────────── REGISTRY ────────────────────────────

LAYOUT_RENDERERS = {
    "hero_cover":             render_hero_cover,
    "agenda_cards":           render_agenda_cards,
    "summary_dashboard":      render_summary_dashboard,
    "kpi_dashboard":          render_kpi_dashboard,
    "dashboard_grid":         render_dashboard_grid,
    "performance_dashboard":  render_performance_dashboard,
    "split_metrics_chart":    render_split_metrics_chart,
    "insight_dashboard":      render_insight_dashboard,
    "three_column_status":    render_three_column_status,
    "donut_insights":         render_donut_insights,
    "risk_dashboard":         render_risk_dashboard,
    "opportunity_dashboard":  render_opportunity_dashboard,
    "action_tracker":         render_action_tracker,
    # Roadmap/timeline visuals are intentionally disabled; any accidental
    # roadmap layout renders as an action tracker instead.
    "recommendation_roadmap": render_action_tracker,
    "roadmap":                render_action_tracker,
    "timeline":               render_action_tracker,
    "comparison_dashboard":   render_comparison_dashboard,
    "before_after":           render_before_after,
    "closing_slide":          render_closing_slide,
}