from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def rgb(hex_color):
    hex_color = str(hex_color or "111827").replace("#", "")
    if len(hex_color) != 6:
        hex_color = "111827"
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def fit_text(text, max_chars=200):
    """Clamp text without adding ellipsis markers.

    PowerPoint cards must not bleed outside their boxes, but the generated deck
    should not show artificial "..." / "…" truncation marks.  When a hard
    clamp is needed, cut at a word boundary and rely on textbox auto-fit/wrap.
    """
    text = str(text or "").strip().replace("…", "")
    text = __import__("re").sub(r"\.{3,}", "", text)
    if len(text) <= max_chars:
        return text
    clipped = text[:max_chars].rstrip()
    last_space = clipped.rfind(" ")
    if last_space > max_chars // 2:
        clipped = clipped[:last_space].rstrip()
    return clipped.rstrip(" ,;:-")


def _set_run_style(run, font_size, bold, color):
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_textbox(slide, text, x, y, w, h, font_size=16, bold=False, color="111827",
                align=PP_ALIGN.LEFT):
    """
    Add a textbox. Text shrinks to fit the box (no overflow).
    Uses NONE auto-size so the box stays at the given size and
    the font shrinks via txBody normAutofit instead.
    """
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = box.text_frame
    frame.word_wrap = True

    # Use NONE so box doesn't resize — we shrink text instead
    frame.auto_size = None  # MSO_AUTO_SIZE.NONE

    # Apply normAutofit (shrink-text-on-overflow) via XML
    txBody = frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        # Remove any existing spAutoFit / noAutofit
        for tag in ["a:spAutoFit", "a:noAutofit", "a:normAutofit"]:
            old = bodyPr.find(qn(tag))
            if old is not None:
                bodyPr.remove(old)
        bodyPr.append(etree.SubElement(bodyPr, qn("a:normAutofit")))

    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text or "")
    _set_run_style(run, font_size, bold, color)

    return box


def add_multiline_textbox(slide, text, x, y, w, h, font_size=13, bold=False, color="111827"):
    """
    Add a textbox that splits on newlines into separate paragraphs,
    so each bullet appears on its own line.
    """
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = None

    txBody = frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for tag in ["a:spAutoFit", "a:noAutofit", "a:normAutofit"]:
            old = bodyPr.find(qn(tag))
            if old is not None:
                bodyPr.remove(old)
        bodyPr.append(etree.SubElement(bodyPr, qn("a:normAutofit")))

    lines = str(text or "").split("\n")
    frame.clear()
    for idx, line in enumerate(lines):
        if idx == 0:
            para = frame.paragraphs[0]
        else:
            para = frame.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line.strip()
        _set_run_style(run, font_size, bold, color)

    return box


def add_title(slide, title, theme, x=0.55, y=0.22, w=12.3, h=0.75):
    colors = theme["colors"]
    return add_textbox(
        slide, title, x, y, w, h,
        font_size=26, bold=True, color=colors["text"]
    )


def add_subtitle(slide, text, theme, x=0.55, y=0.98, w=12.0, h=0.48):
    colors = theme["colors"]
    return add_textbox(
        slide, text, x, y, w, h,
        font_size=13, bold=False, color=colors["subtext"]
    )


def add_background(slide, theme):
    colors = theme["colors"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(colors["background"])


def style_card_shape(shape, theme):
    shapes_cfg = theme.get("shapes", {}) if isinstance(theme, dict) else {}

    radius = shapes_cfg.get("card_radius")
    if radius is not None:
        try:
            shape.adjustments[0] = max(0.0, min(float(radius), 0.5))
        except Exception:
            pass
    return shape


def add_card(slide, title, body, theme, x, y, w, h):
    """
    Render a card using the shape's OWN text_frame so text is always
    clipped to the shape boundary — no floating textboxes that overflow.
    Title is the first paragraph (bold), then a blank line, then body lines.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor as _RGB
    from lxml import etree

    colors = theme["colors"]

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colors["surface"])
    shape.line.color.rgb = rgb(colors["primary"])
    shape.line.width = Emu(12700)
    style_card_shape(shape, theme)

    # ── use the shape's built-in text_frame (clips to shape) ──
    tf = shape.text_frame
    tf.word_wrap = True

    # Set internal padding (top/left/right/bottom) via bodyPr attributes
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        _pad = int(Inches(0.14))
        bodyPr.set("lIns", str(_pad))
        bodyPr.set("rIns", str(_pad))
        bodyPr.set("tIns", str(_pad))
        bodyPr.set("bIns", str(int(Inches(0.10))))
        bodyPr.set("anchor", "t")  # top-align text
        # Use normAutofit so text shrinks to fit rather than overflowing
        for tag in ["a:spAutoFit", "a:noAutofit", "a:normAutofit"]:
            el = bodyPr.find(qn(tag))
            if el is not None:
                bodyPr.remove(el)
        bodyPr.append(etree.SubElement(bodyPr, qn("a:normAutofit")))

    # Scale font sizes to card height
    title_fs = Pt(14 if h >= 3.0 else 13)
    body_fs  = Pt(12 if h >= 3.0 else (11 if h >= 1.8 else 10))

    title_str = fit_text(str(title or ""), 80)
    body_str  = fit_text(str(body or ""), 500)

    # ── Paragraph 1: title (bold) ──
    p0 = tf.paragraphs[0]
    p0.clear()
    run0 = p0.add_run()
    run0.text = title_str
    run0.font.bold = True
    run0.font.size = title_fs
    run0.font.color.rgb = _RGB.from_string(colors["text"])
    p0.space_after = Pt(4)

    # ── Remaining paragraphs: body lines ──
    lines = [l for l in body_str.split("\n") if l.strip()]
    for line in lines:
        para = tf.add_paragraph()
        para.clear()
        run = para.add_run()
        run.text = line
        run.font.bold = False
        run.font.size = body_fs
        run.font.color.rgb = _RGB.from_string(colors["subtext"])
        para.space_before = Pt(2)
        para.space_after  = Pt(2)

    return shape


def add_kpi_card(slide, label, value, note, theme, x, y, w, h):
    colors = theme["colors"]

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colors["surface"])
    shape.line.color.rgb = rgb(colors["secondary"])
    shape.line.width = Emu(12700)
    style_card_shape(shape, theme)
    shape.text = ""

    # KPI cards are reused both as large dashboard tiles and as compact
    # top-strip tiles.  Compact tiles cannot safely carry long paragraph notes;
    # aggressively shorten text by available area so it stays inside the shape
    # instead of bleeding into lower boxes.
    compact = h <= 1.55
    area = max(w * h, 0.1)
    label_limit = 42 if compact else 80
    value_limit = 28 if compact else 40
    note_limit = int(max(48, min(115 if compact else 260, area * (18 if compact else 34))))

    label_str = fit_text(str(label or ""), label_limit)
    raw_value = str(value or "")
    raw_note = str(note or "")
    # If a KPI value is a long list (e.g., segment ACV), push most detail into the note
    # and keep the large value line short so it cannot bleed outside the box.
    if len(raw_value) > value_limit or "," in raw_value:
        # Long KPI values such as segment-level ACV should not be shown as a
        # giant clipped number. Put the detail in the wrapped note area and use
        # a compact value label with no ellipsis.
        value_str = "By segment" if "," in raw_value else "See detail"
        raw_note = (raw_note + " " + raw_value).strip() if raw_note else raw_value
    else:
        value_str = fit_text(raw_value, value_limit)
    note_str  = fit_text(raw_note, note_limit)

    if compact:
        pad_x = 0.16
        top_pad = 0.10
        label_h = 0.24
        value_h = 0.38 if len(value_str) <= 12 else 0.34
        note_y = y + top_pad + label_h + value_h + 0.08
        note_h = max(0.18, y + h - 0.13 - note_y)
        label_font = 8.5
        value_font = 18 if len(value_str) <= 12 else 15
        note_font = 8
    else:
        pad_x = 0.15
        top_pad = 0.10
        label_h = max(0.26, h * 0.14)
        value_h = max(0.42, h * 0.35)
        note_y = y + label_h + 0.10 + value_h + 0.08
        note_h = max(0.20, (y + h - 0.12) - note_y)
        label_font = 10
        value_font = (28 if h >= 2.0 else 22) if len(value_str) <= 12 else (20 if h >= 2.0 else 16)
        note_font = 10

    inner_w = max(0.2, w - (pad_x * 2))

    add_textbox(
        slide, label_str,
        x + pad_x, y + top_pad, inner_w, label_h,
        font_size=label_font, bold=False, color=colors["subtext"]
    )

    add_textbox(
        slide, value_str,
        x + pad_x, y + top_pad + label_h + 0.02, inner_w, value_h,
        font_size=value_font, bold=True, color=colors["accent"]
    )

    if note_str and note_h > 0.16:
        add_textbox(
            slide, note_str,
            x + pad_x, note_y, inner_w, note_h,
            font_size=note_font, bold=False, color=colors["subtext"]
        )
    return shape


def add_section_label(slide, text, theme, x, y, w=2.4, h=0.30):
    colors = theme["colors"]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colors["accent"])
    shape.line.color.rgb = rgb(colors["accent"])

    frame = shape.text_frame
    frame.word_wrap = False
    frame.auto_size = None
    frame.clear()
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = fit_text(str(text or ""), 35)
    _set_run_style(run, 9, True, "FFFFFF")
    return shape


def add_bullet_list(slide, items, theme, x, y, w, h, font_size=13):
    colors = theme["colors"]
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = None

    txBody = frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for tag in ["a:spAutoFit", "a:noAutofit", "a:normAutofit"]:
            old = bodyPr.find(qn(tag))
            if old is not None:
                bodyPr.remove(old)
        bodyPr.append(etree.SubElement(bodyPr, qn("a:normAutofit")))

    frame.clear()
    for idx, item in enumerate(items or []):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = "• " + fit_text(str(item), 200)
        _set_run_style(run, font_size, False, colors["subtext"])
    return box


def add_metric_strip(slide, cards, theme, x=0.55, y=1.38, w=12.2, h=1.2):
    if not cards:
        return
    gap = 0.16
    count = min(len(cards), 5)
    card_w = (w - gap * (count - 1)) / count

    for i, card in enumerate(cards[:count]):
        cx = x + i * (card_w + gap)
        label = card.get("name") or card.get("label") or "Metric"
        value = card.get("value", "")
        note = card.get("interpretation") or card.get("note") or ""
        add_kpi_card(slide, label, value, note, theme, cx, y, card_w, h)


def add_progress_bar(slide, label, value, theme, x, y, w, h=0.26):
    colors = theme["colors"]
    try:
        numeric = float(str(value).replace("%", "").strip())
    except Exception:
        numeric = 0
    numeric = max(0, min(numeric, 100))

    add_textbox(slide, label, x, y - 0.26, w, 0.24,
                font_size=10, bold=True, color=colors["text"])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(colors["surface"])
    bg.line.color.rgb = rgb(colors["surface"])

    fill_w = max(w * numeric / 100, 0.02)
    fg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(fill_w), Inches(h)
    )
    fg.fill.solid()
    fg.fill.fore_color.rgb = rgb(colors["accent"])
    fg.line.color.rgb = rgb(colors["accent"])

    add_textbox(slide, f"{numeric:.0f}%", x + w - 0.55, y - 0.03,
                0.5, 0.24, font_size=9, bold=True, color=colors["text"])


def add_status_card(slide, label, title, note, theme, x, y, w, h, status="neutral"):
    colors = theme["colors"]
    status_color_map = {
        "risk": colors["danger"],
        "success": colors["success"],
        "action": colors["accent"],
    }
    status_color = status_color_map.get(status, colors["primary"])

    label_str = fit_text(str(label or ""), 45)
    title_str = fit_text(str(title or ""), 125)
    note_str  = fit_text(str(note or ""), 190)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colors["surface"])
    shape.line.color.rgb = rgb(status_color)
    shape.line.width = Emu(19050)  # 1.5pt
    style_card_shape(shape, theme)
    shape.text = ""

    label_w = min(len(label_str) * 0.095 + 0.4, w - 0.3)
    add_section_label(slide, label_str, theme, x + 0.16, y + 0.13, w=label_w, h=0.26)

    title_h = min(h * 0.38, 0.5)
    add_textbox(
        slide, title_str,
        x + 0.16, y + 0.46, w - 0.32, title_h,
        font_size=12 if len(title_str) > 85 else 13, bold=True, color=colors["text"]
    )
    if note_str:
        add_textbox(
            slide, note_str,
            x + 0.16, y + 0.46 + title_h + 0.04, w - 0.32,
            max(h - 0.46 - title_h - 0.1, 0.2),
            font_size=10, bold=False, color=colors["subtext"]
        )
    return shape


def add_header_band(slide, theme, height=1.45):
    colors = theme["colors"]
    shapes_cfg = theme.get("shapes", {}) if isinstance(theme, dict) else {}
    if not shapes_cfg.get("header_band"):
        return None

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(height)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(colors["surface"])
    band.line.fill.background()

    spTree = slide.shapes._spTree
    spTree.remove(band._element)
    spTree.insert(2, band._element)
    return band


def add_footer(slide, theme, page=None):
    colors = theme["colors"]
    text = "Executive Insight Generator"
    if page:
        text = f"{text}  ·  {page}"
    add_textbox(
        slide, text,
        0.55, 7.12, 12.2, 0.22,
        font_size=8, bold=False, color=colors["subtext"]
    )


def add_divider_line(slide, theme, x, y, w, accent=False):
    colors = theme["colors"]
    color = colors["accent"] if accent else colors["subtext"]
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(color)
    line.line.fill.background()
    return line
