"""
layout_utils.py — shared safe-placement helpers
Imported by layout_library.py to prevent any shape from overflowing the slide.
"""
from pptx.util import Inches, Emu

SLIDE_W = 13.33   # inches
SLIDE_H = 7.50    # inches
FOOTER_TOP = 7.12  # footer starts here — content must stay above this
SAFE_BOTTOM = FOOTER_TOP - 0.06   # 7.06" — last pixel before footer


def safe_inches(val, lo=0.0, hi=None):
    """Clamp a value (in inches) between lo and hi."""
    if hi is None:
        hi = SLIDE_W
    return max(lo, min(float(val), hi))


def safe_box(x, y, w, h, margin=0.0):
    """
    Return (x, y, w, h) guaranteed to fit inside the slide content area.
    Content area: x in [0, SLIDE_W], y in [0, SAFE_BOTTOM].
    """
    x = safe_inches(x, 0.0, SLIDE_W)
    y = safe_inches(y, 0.0, SAFE_BOTTOM)
    # Width: must not exceed slide right edge
    max_w = SLIDE_W - x - margin
    w = safe_inches(w, 0.01, max_w)
    # Height: must not push bottom below footer
    max_h = SAFE_BOTTOM - y - margin
    h = safe_inches(h, 0.01, max_h)
    return x, y, w, h


def add_picture_safe(slide, path, x, y, w, h=None, reserved_top=1.55, reserved_bottom=None):
    """
    Add a picture and auto-scale its height so it never overflows.
    - reserved_top: y coordinate of the top of the content zone (after title/subtitle/divider)
    - w: desired width in inches
    - h: optional explicit height in inches; use this for equal chart cells
    Returns True on success, False on failure.
    """
    if reserved_bottom is None:
        reserved_bottom = SAFE_BOTTOM

    max_h = reserved_bottom - y
    if max_h < 0.5:
        return False

    try:
        from pptx.util import Inches as I
        if h is not None:
            h = min(float(h), max_h)
            if h < 0.5:
                return False
            slide.shapes.add_picture(path, I(x), I(y), width=I(w), height=I(h))
            return True

        pic = slide.shapes.add_picture(path, I(x), I(y), width=I(w))
        # Check if height overflows and scale down proportionally
        pic_h = pic.height / 914400  # EMU → inches
        if pic_h > max_h:
            scale = max_h / pic_h
            pic.width  = int(pic.width  * scale)
            pic.height = int(pic.height * scale)
        return True
    except Exception as e:
        print(f"[layout_utils] picture error: {e}")
        return False


def content_height(reserved_top=1.55):
    """Usable slide height between reserved_top and the footer."""
    return SAFE_BOTTOM - reserved_top
