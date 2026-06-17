import os
import re
import pdfplumber
import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


# ─── cell / number helpers ───────────────────────────────────────

def clean_cell(value):
    if value is None:
        return ""
    return str(value).strip()


def parse_number(value):
    if value is None:
        return None
    text = str(value)
    # Handle M/B/K suffix BEFORE stripping them
    multiplier = 1
    stripped = text.replace("$", "").replace(",", "").replace("%", "").strip()
    # Support both compact units (18.7B) and prose units (18.7 billion).
    unit_match = re.search(r"\b(billion|million|thousand)\b$", stripped, re.IGNORECASE)
    if unit_match:
        unit = unit_match.group(1).lower()
        multiplier = {"billion": 1e9, "million": 1e6, "thousand": 1e3}[unit]
        stripped = re.sub(r"\s*\b(billion|million|thousand)\b$", "", stripped, flags=re.IGNORECASE).strip()
    elif stripped.upper().endswith("B"):
        multiplier = 1e9
        stripped = stripped[:-1].strip()
    elif stripped.upper().endswith("M"):
        multiplier = 1e6
        stripped = stripped[:-1].strip()
    elif stripped.upper().endswith("K"):
        multiplier = 1e3
        stripped = stripped[:-1].strip()
    try:
        return float(stripped) * multiplier
    except Exception:
        return None


def extract_numbers_from_text(text):
    pattern = r"[-+]?\$?\d[\d,]*(?:\.\d+)?[BMK]?%?"
    matches = re.findall(pattern, text)

    numbers = []
    for match in matches:
        value = parse_number(match)
        if value is not None:
            numbers.append({"raw": match, "value": value})
    return numbers


def dataframe_to_table(df, table_name="Table"):
    df = df.fillna("")
    headers = [str(col).strip() for col in df.columns]
    rows = [[clean_cell(v) for v in row.tolist()] for _, row in df.iterrows()]
    return {"table_name": table_name, "headers": headers, "rows": rows}


def extract_docx(path):
    doc = Document(path)
    text_parts, tables = [], []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    for index, table in enumerate(doc.tables):
        rows = [[clean_cell(cell.text) for cell in row.cells] for row in table.rows]
        if len(rows) > 1:
            tables.append({
                "table_name": f"DOCX Table {index + 1}",
                "headers": rows[0],
                "rows": rows[1:]
            })
    return "\n".join(text_parts), tables


def extract_pdf(path):
    text_parts, tables = [], []
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"--- Page {page_num+1} ---\n{page_text}")
            for idx, table in enumerate(page.extract_tables() or []):
                if len(table) < 2:
                    continue
                headers = [clean_cell(x) for x in table[0]]
                rows    = [[clean_cell(cell) for cell in row] for row in table[1:]]
                tables.append({
                    "table_name": f"PDF Page {page_num+1} Table {idx+1}",
                    "headers": headers,
                    "rows": rows
                })
    text = "\n\n".join(text_parts)
    if not tables:
        tables = detect_text_tables(text)

    return text, tables

def extract_pptx(path):
    prs = Presentation(path)
    text_parts, tables = [], []
    for slide_index, slide in enumerate(prs.slides):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            if shape.has_table:
                table = shape.table
                rows = [[clean_cell(cell.text) for cell in row.cells] for row in table.rows]
                if len(rows) > 1:
                    tables.append({
                        "table_name": f"PPTX Slide {slide_index+1} Table",
                        "headers": rows[0],
                        "rows": rows[1:]
                    })
        if slide_text:
            text_parts.append(f"--- Slide {slide_index+1} ---\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts), tables


def extract_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return text, []


def extract_csv(path):
    df = pd.read_csv(path, skipinitialspace=True, on_bad_lines="skip")
    table = dataframe_to_table(df, "CSV Table")
    return df.to_string(index=False), [table]


def extract_excel(path):
    ext = os.path.splitext(path)[1].lower()
    engine = "xlrd" if ext == ".xls" else "openpyxl"
    xl = pd.ExcelFile(path, engine=engine)
    text_parts, tables = [], []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet).fillna("")
        text_parts.append(f"Sheet: {sheet}\n{df.to_string(index=False)}")
        tables.append(dataframe_to_table(df, f"Sheet: {sheet}"))
    return "\n\n".join(text_parts), tables


# ─── KPI / metric extraction ─────────────────────────────────────

# Explicit KPI patterns extracted from prose like earnings call transcripts
_PROSE_KPI_PATTERNS = [
    # "revenue of $18.7 billion"
    (r"(?:total\s+)?revenue[s]?\s+(?:of|were|was|reached|grew to)\s+\$?([\d,.]+)\s*(billion|million|B|M)?",
     "Revenue"),
    # "bookings of $20.9 billion"
    (r"(?:new\s+)?bookings\s+(?:of|were|was)\s+\$?([\d,.]+)\s*(billion|million|B|M)?",
     "Bookings"),
    # "adjusted EPS of $3.94"
    (r"(?:adjusted\s+)?(?:diluted\s+)?EPS\s+(?:of|were|was|in the quarter of)\s+\$?([\d,.]+)",
     "Adjusted EPS"),
    # "adjusted operating margin was 17%"
    (r"(?:adjusted\s+)?operating\s+margin\s+(?:was|of|is)\s+([\d.]+)%",
     "Operating Margin"),
    # "free cash flow of $1.5 billion"
    (r"free\s+cash\s+flow\s+(?:of|was|were)\s+\$?([\d,.]+)\s*(billion|million|B|M)?",
     "Free Cash Flow"),
    # "grew 5% in local currency"
    (r"(?:revenue[s]?|growth)\s+(?:grew|growing|grow)\s+([\d.]+)%\s+in\s+local\s+currency",
     "Revenue Growth (LC)"),
    # "EPS growth of 10%"
    (r"EPS\s+growth\s+of\s+([\d.]+)%",
     "EPS Growth"),
    # "advanced AI bookings ... $2.2 billion"
    (r"advanced\s+AI\s+bookings\s+(?:this quarter were|of)\s+\$?([\d,.]+)\s*(billion|million|B|M)?",
     "Advanced AI Bookings"),
    # "returned $3.3 billion to shareholders"
    (r"returned\s+\$?([\d,.]+)\s*(billion|million|B|M)?\s+to\s+shareholders",
     "Returned to Shareholders"),
    # "30 basis points"
    (r"(?:margin\s+)?expanded?\s+by\s+([\d.]+)\s+basis\s+points",
     "Margin Expansion (bps)"),
    # "book-to-bill of 1.1"
    (r"book-to-bill\s+of\s+([\d.]+)",
     "Book-to-Bill"),
    # "gross margin ... 33.1%"
    (r"(?:gross\s+)?margin\s+for the quarter was\s+([\d.]+)%",
     "Gross Margin"),
    # "tax rate ... 23.9%"
    (r"(?:effective\s+)?tax[- ]rate\s+(?:for the quarter was|of)\s+([\d.]+)%",
     "Effective Tax Rate"),
    # "headcount ... 784,000"
    (r"(?:nearly|approximately)\s+([\d,]+)\s+(?:people|professionals|employees)",
     "Headcount"),
    # "cash balance ... $9.6 billion"
    (r"cash\s+balance\s+(?:at|was|of)\s+\$?([\d,.]+)\s*(billion|million|B|M)?",
     "Cash Balance"),
    # "repurchased ... 9.5 million shares"
    (r"repurchased\s+(?:or\s+redeemed\s+)?([\d,.]+)\s*million\s+shares",
     "Shares Repurchased (M)"),
    # "dividend of $1.63 per share"
    (r"(?:cash\s+)?dividend\s+of\s+\$?([\d.]+)\s+per\s+share",
     "Dividend Per Share"),
]

_MULTIPLIERS = {"billion": 1e9, "million": 1e6, "b": 1e9, "m": 1e6}

def _format_value(raw_num: str, unit: str = "") -> str:
    """Format extracted number+unit back into a readable metric value string."""
    try:
        v = float(raw_num.replace(",", ""))
    except Exception:
        return raw_num

    unit = (unit or "").lower().strip()
    mult = _MULTIPLIERS.get(unit, 1)

    if mult >= 1e9:
        return f"${v:.1f}B"
    elif mult >= 1e6:
        return f"${v:.1f}M"
    elif "%" in raw_num or "margin" in unit or "growth" in unit or "rate" in unit or "bps" in unit:
        return f"{v:.1f}%"
    else:
        return raw_num


def extract_metric_cards_from_prose(text: str, max_cards: int = 30) -> list:
    """
    Extract KPI cards from prose/transcript text using financial sentence patterns.
    Returns [{label, value}] dicts.
    """
    cards = []
    seen  = set()

    for pattern, label in _PROSE_KPI_PATTERNS:
        for m in re.finditer(pattern, text, re.IGNORECASE):
            groups = m.groups()
            raw_num = groups[0].replace(",", "") if groups else ""
            unit    = groups[1] if len(groups) > 1 and groups[1] else ""

            try:
                float(raw_num)
            except Exception:
                continue

            value = _format_value(raw_num, unit)
            key   = (label.lower(), value.lower())
            if key in seen:
                continue
            seen.add(key)
            cards.append({"label": label, "value": value})
            if len(cards) >= max_cards:
                return cards

    return cards


def extract_metric_cards(text, max_cards=60):
    """
    Combined extractor: tries structured Label: Value patterns first,
    then falls back to prose sentence patterns for transcripts / reports.
    """
    cards = []
    seen  = set()

    # --- Pattern 1: "Label: $Value" or "Label    Value" (structured docs) ---
    for line in text.splitlines():
        line = line.strip()
        if len(line) < 5 or len(line) > 140:
            continue

        # "Label: Value"
        m = re.match(
            r"^(.{3,70}?):\s*([$₹€£]?\d[\d,]*(?:\.\d+)?\s*[%A-Za-z/µμ₹$€£.-]*)$",
            line
        )
        if m:
            label = m.group(1).strip()
            value = m.group(2).strip()
            if parse_number(value) is not None:
                key = (label.lower(), value.lower())
                if key not in seen:
                    cards.append({"label": label, "value": value})
                    seen.add(key)
                if len(cards) >= max_cards:
                    break
            continue

        # "Label    Value   (whitespace-separated)"
        parts = re.split(r"\s{2,}", line)
        if len(parts) >= 2:
            label = parts[0].strip()
            value = parts[1].strip()
            if parse_number(value) is not None:
                key = (label.lower(), value.lower())
                if key not in seen:
                    cards.append({"label": label, "value": value})
                    seen.add(key)
                if len(cards) >= max_cards:
                    break

    # --- Fallback: prose sentence patterns (earnings calls, transcripts, reports) ---
    if len(cards) < 5:
        prose_cards = extract_metric_cards_from_prose(text, max_cards=max_cards - len(cards))
        for c in prose_cards:
            key = (c["label"].lower(), c["value"].lower())
            if key not in seen:
                cards.append(c)
                seen.add(key)

    return cards[:max_cards]


# ─── chart candidate builder ─────────────────────────────────────

def _looks_time_label(label: str) -> bool:
    label_l = label.lower()
    return bool(re.search(r"\b(jan|feb|mar|apr|may|jun|jul|aug|sep|sept|oct|nov|dec|q[1-4]|20\d{2}|19\d{2})\b", label_l))


def _clean_chart_label(label: str) -> str:
    label = re.sub(r"^[,;:\-–—\s]+|[,;:\-–—\s]+$", "", str(label or "").strip())
    return re.sub(r"\s+", " ", label)



def _is_research_like_text(text: str) -> bool:
    """Heuristic guard: research papers contain many citations and method words.

    For these documents, random years/page numbers/citations should not become
    executive charts or metric cards.  The RAG/reasoning layer can still use the
    text, but chart extraction must be conservative.
    """
    low = str(text or "").lower()
    research_hits = sum(1 for t in (
        "abstract", "method", "methodology", "participants", "sample",
        "research question", "data collection", "analysis", "literature",
        "references", "citation", "journal", "hypothesis", "study", "findings",
    ) if t in low)
    business_hits = sum(1 for t in (
        "revenue", "bookings", "gross margin", "operating margin", "ebitda",
        "cash flow", "sales pipeline", "contract value", "inventory", "customer churn",
    ) if t in low)
    return research_hits >= 4 and research_hits > business_hits


def _is_bad_chart_line(line: str, title: str = "") -> bool:
    """Reject citation/debug lines that look numeric but are not chartable."""
    low = f"{title} {line}".lower()
    if any(x in low for x in (
        "doi", "isbn", "http", "www.", "retrieved from", "references", "bibliography",
        "cavus", "ibrahim", "journal", "vol.", "no.", "pp.", "page ", " p.",
        "pdf page", "table row", "line-level metrics", "document section",
    )):
        return True
    # Citation-like years without a real metric label should not be charted.
    if re.search(r"\b(19|20)\d{2}\b", low) and not any(k in low for k in (
        "revenue", "sales", "margin", "growth", "expense", "cash", "arr", "nrr",
        "churn", "conversion", "participants", "sample size", "accuracy", "score",
    )):
        return True
    return False


def _should_use_pie_chart(title: str, label_col: str = "", value_col: str = "", data=None) -> bool:
    """Only use pie/donut for true composition/share charts.

    KPI scorecards, growth rates, margins, variances, time periods, and mixed
    metric lists are much clearer as bar/line charts and can break when values
    are negative.  This generic rule prevents noisy pies across all domains.
    """
    data = data or []
    text = f"{title} {label_col} {value_col}".lower()
    try:
        vals = [float(d.get("value")) for d in data]
    except Exception:
        return False
    if len(vals) < 2 or len(vals) > 6 or any(v < 0 for v in vals) or sum(vals) <= 0:
        return False
    blocked = (
        "growth", "margin", "rate", "variance", "actual", "plan", "fy24", "fy25",
        "fy26", "metric", "kpi", "scorecard", "eps", "tax", "days", "cycle",
        "line-level", "cash conversion", "churn", "nps", "ebitda", "operating margin",
    )
    if any(b in text for b in blocked):
        return False
    allowed = (
        "mix", "share", "distribution", "composition", "portfolio allocation",
    )
    return any(a in text for a in allowed)


def build_chart_candidates_from_text_lines(text: str, max_candidates: int = 12) -> list:
    """Build chart candidates from prose/line-level numeric statements.

    This catches patterns such as:
    - "Revenue by month: Jan $10M, Feb $12M, Mar $14M"
    - "North America 42%  Europe 31%  APAC 27%"
    - "Q1 12.4, Q2 13.8, Q3 14.1"

    It complements table extraction instead of requiring a formal table.
    """
    candidates, seen = [], set()
    number_token = r"[$₹€£]?\d[\d,]*(?:\.\d+)?\s*(?:%|B|M|K|billion|million|thousand)?"
    pair_re = re.compile(rf"([A-Za-z][A-Za-z0-9/&().% '\-]{{1,38}}?)\s+({number_token})", re.IGNORECASE)

    for raw_line in text.splitlines():
        line = re.sub(r"\s+", " ", raw_line.strip())
        if len(line) < 10 or len(line) > 220:
            continue
        if sum(ch.isdigit() for ch in line) < 2:
            continue
        if _is_bad_chart_line(line):
            continue

        prefix = ""
        body = line
        if ":" in line:
            prefix, body = line.split(":", 1)
            prefix = _clean_chart_label(prefix)
        else:
            line_l = line.lower()
            if not any(k in line_l for k in (" by ", "q1", "q2", "q3", "q4", "jan", "feb", "mar", "apr", "may", "jun", "jul", "aug", "sep", "oct", "nov", "dec", "share", "mix", "rate", "margin", "revenue", "expense", "cash")):
                continue

        pairs = []
        for m in pair_re.finditer(body):
            label = _clean_chart_label(m.group(1))
            value_raw = m.group(2).strip()
            if not label or len(label.split()) > 5:
                continue
            if label.lower() in {"and", "or", "to", "from", "with", "by", "of", "in", "at"}:
                continue
            value = parse_number(value_raw)
            if value is None:
                continue
            pairs.append({"label": label, "value": value})

        # De-duplicate while preserving order.
        cleaned, label_seen = [], set()
        for p in pairs:
            key = p["label"].lower()
            if key not in label_seen:
                cleaned.append(p)
                label_seen.add(key)

        if len(cleaned) < 2:
            continue

        title = prefix if 3 <= len(prefix) <= 70 else "Line-Level Metrics"
        if title.lower() in {"line-level metrics", "metrics from document lines"}:
            # Generic numeric fragments are usually citations or mixed metrics;
            # they should not become visible charts.
            continue
        if title.lower() in {"key metrics", "metrics", "summary"}:
            title = "Metrics from Document Lines"
        if _is_bad_chart_line(line, title):
            continue
        key = (title.lower(), tuple((p["label"].lower(), p["value"]) for p in cleaned[:8]))
        if key in seen:
            continue
        seen.add(key)

        ctype = "line" if any(_looks_time_label(p["label"]) for p in cleaned) else "bar"
        if _should_use_pie_chart(title, data=cleaned):
            ctype = "pie"

        candidates.append({
            "title": title,
            "chart_type": ctype,
            "data": cleaned[:12],
        })
        if len(candidates) >= max_candidates:
            break

    return candidates


def build_chart_candidates_from_tables(tables, metric_cards=None):
    """
    Build chart candidates from:
    1. Structured tables (original logic)
    2. Prose-extracted metric cards (new — for transcripts with no tables)
    """
    candidates = []
    seen_titles = set()

    # --- From tables ---
    for table in tables:
        headers = table.get("headers", [])
        rows    = table.get("rows", [])
        name    = table.get("table_name", "Table")

        if len(headers) < 2 or len(rows) < 2:
            continue

        label_col = headers[0]
        for value_col in headers[1:]:
            data = []
            for row in rows:
                if len(row) < 2:
                    continue
                label = clean_cell(row[0])
                value = parse_number(row[headers.index(value_col)] if value_col in headers else row[1])
                if label and value is not None:
                    data.append({"label": label, "value": value})

            if len(data) >= 2:
                title = f"{value_col} by {label_col}" if label_col else value_col
                key   = title.lower()
                if key not in seen_titles:
                    # Decide chart type.  Default to bar/line; only use pie for true share/mix charts.
                    ctype = "bar"
                    if any(h.lower() in ("month", "quarter", "year", "date", "period")
                           for h in [label_col]):
                        ctype = "line"
                    elif _should_use_pie_chart(title, label_col, value_col, data):
                        ctype = "pie"
                    candidates.append({
                        "title":      title,
                        "chart_type": ctype,
                        "data":       data
                    })
                    seen_titles.add(key)

    # --- From metric cards (for prose/transcript docs with no tables) ---
    if metric_cards:
        # Group into meaningful chart clusters
        revenue_cards = [c for c in metric_cards
                         if any(k in c["label"].lower()
                                for k in ("revenue", "bookings", "cash", "eps", "income"))]
        margin_cards  = [c for c in metric_cards
                         if any(k in c["label"].lower()
                                for k in ("margin", "growth", "rate", "basis", "tax"))]

        if len(revenue_cards) >= 2:
            data = []
            for c in revenue_cards[:6]:
                v = parse_number(c["value"])
                if v is not None:
                    data.append({"label": c["label"], "value": v})
            if len(data) >= 2:
                candidates.append({
                    "title": "Key Financial Metrics",
                    "chart_type": "bar",
                    "data": data
                })

        if len(margin_cards) >= 2:
            data = []
            for c in margin_cards[:6]:
                v = parse_number(c["value"])
                if v is not None:
                    data.append({"label": c["label"], "value": v})
            if len(data) >= 2:
                candidates.append({
                    "title": "Margin & Growth Metrics",
                    "chart_type": "bar",
                    "data": data
                })

        # Do not create hard-coded mix charts.  If a document has a real
        # consulting/managed-services table, table extraction above will create it.

    return candidates


# ─── section extractor ───────────────────────────────────────────

# Words that appear as headings in transcripts but are not real document sections
_JUNK_SECTION_WORDS = {
    "thank", "thanks", "good", "morning", "happy", "holiday", "operator", "welcome",
    "please", "forward", "looking", "gaap", "non", "sec", "filing", "form", "call",
    "transcript", "quarter", "fiscal", "results", "question", "answer", "participant",
    "apologies", "briefly", "let me", "just", "really", "know", "see", "like",
    "think", "back", "now", "also", "going", "right", "one", "two", "three",
    "well", "so", "we", "our", "the", "this", "that", "with", "from", "into",
    "about", "there", "which", "where", "how", "what", "when", "who", "why",
    "if", "as", "it", "in", "is", "are", "was", "were", "have", "has", "had",
    "be", "been", "being", "do", "does", "did", "will", "would", "could", "should",
    "may", "might", "must", "shall", "can", "need", "ought", "used", "dare",
}

def extract_sections(text):
    """
    Extract meaningful section headings.
    Filters out sentence fragments, filler words, and transcript chatter.
    """
    sections = []
    seen     = set()

    # Original pattern: lines that look like headings (title case, short)
    for line in text.splitlines():
        line = line.strip()
        if not line or len(line) > 80 or len(line) < 5:
            continue

        # Must look like a heading: mostly title-case or ALL CAPS, few words
        words = line.split()
        if len(words) > 8:
            continue

        # Filter junk words
        lower_words = [w.lower().strip(".,;:!?") for w in words]
        if any(w in _JUNK_SECTION_WORDS for w in lower_words):
            continue

        # Must start with capital
        if not line[0].isupper():
            continue

        # Must not be a sentence (ends with . or ?)
        if line.endswith((".","?","!")):
            continue

        key = line.lower()
        if key not in seen:
            sections.append(line)
            seen.add(key)

    # Also add known structural sections from transcripts
    structural = re.findall(
        r"\b(Digital Core|Advanced AI|Managed Services|Consulting|"
        r"Industry X|Accenture Song|Partnership Strategy|Market Outlook|"
        r"Business Outlook|Financial Results|Geographic Markets|"
        r"Capital Allocation|Balance Sheet|Operating Margin|"
        r"Revenue Growth|Bookings|Free Cash Flow)\b",
        text,
        re.IGNORECASE
    )
    for s in structural:
        key = s.lower()
        if key not in seen:
            sections.append(s)
            seen.add(key)

    return list(dict.fromkeys(sections))  # preserve order, deduplicate


# ─── text table detector ─────────────────────────────────────────

def detect_text_tables(text):
    tables = []
    rows   = []

    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        parts = re.split(r"\s{2,}", line)
        if len(parts) >= 3:
            rows.append(parts)

    if len(rows) > 3:
        max_cols = max(len(r) for r in rows)
        headers  = [f"Column {i+1}" for i in range(max_cols)]
        normalized = []
        for r in rows:
            while len(r) < max_cols:
                r.append("")
            normalized.append(r)
        tables.append({
            "table_name": "Detected Text Table",
            "headers":    headers,
            "rows":       normalized
        })

    return tables


# ─── single-column table normalizer ─────────────────────────────

def normalize_single_column_tables(tables):
    def is_number_token(t):
        return re.fullmatch(r"[-+]?\d[\d,]*(?:\.\d+)?%?", t) is not None

    def is_numeric_connector(t):
        return t in ["-", "–", "—", "to", "±", "+/-"]

    def split_row_text(text):
        tokens = text.split()
        if len(tokens) < 2:
            return []
        chunks, text_buffer, i = [], [], 0
        while i < len(tokens):
            token = tokens[i]
            if is_number_token(token):
                if text_buffer:
                    chunks.append(" ".join(text_buffer).strip())
                    text_buffer = []
                numeric_chunk = [token]
                if i + 1 < len(tokens):
                    nt = tokens[i + 1]
                    if not is_number_token(nt) and not is_numeric_connector(nt) and len(nt) <= 12:
                        numeric_chunk.append(nt)
                        i += 1
                if i + 2 < len(tokens):
                    connector = tokens[i + 1]
                    next_num  = tokens[i + 2]
                    if is_numeric_connector(connector) and is_number_token(next_num):
                        numeric_chunk.extend([connector, next_num])
                        i += 2
                        if i + 1 < len(tokens):
                            ar = tokens[i + 1]
                            if not is_number_token(ar) and not is_numeric_connector(ar) and len(ar) <= 12:
                                numeric_chunk.append(ar)
                                i += 1
                chunks.append(" ".join(numeric_chunk).strip())
            else:
                text_buffer.append(token)
            i += 1
        if text_buffer:
            chunks.append(" ".join(text_buffer).strip())
        return [c for c in chunks if c]

    normalized_tables = []
    for table in tables:
        headers = table.get("headers", [])
        rows    = table.get("rows", [])
        if not rows:
            normalized_tables.append(table)
            continue
        is_single_column = len(headers) == 1 or all(len(r) == 1 for r in rows)
        if not is_single_column:
            normalized_tables.append(table)
            continue
        structured_rows = []
        for row in rows:
            if not row:
                continue
            text = clean_cell(row[0])
            if not text:
                continue
            parts = split_row_text(text)
            if len(parts) >= 2:
                structured_rows.append(parts)
        if len(structured_rows) < 2:
            normalized_tables.append(table)
            continue
        max_cols = max(len(r) for r in structured_rows)
        padded   = []
        for r in structured_rows:
            while len(r) < max_cols:
                r.append("")
            padded.append(r)
        normalized_tables.append({
            "table_name": table.get("table_name", "Normalized Table"),
            "headers":    [f"Column {i+1}" for i in range(max_cols)],
            "rows":       padded
        })
    return normalized_tables



# ─── table-derived metric cards ──────────────────────────────────

def extract_metric_cards_from_tables(tables, max_cards=20):
    """Promote KPI scorecard rows to metric cards.

    This is generic: if a table has a metric/KPI column and an actual/current/value
    column, the row becomes a KPI card.  It avoids relying only on prose regexes.
    """
    cards, seen = [], set()
    preferred_value_headers = (
        "actual", "current", "fy25 actual", "fy26", "november 30, 2025",
        "value", "amount", "revenues", "revenue", "result", "score",
    )
    label_headers = ("kpi", "metric", "measure", "indicator", "item", "name")
    for table in tables or []:
        headers = [str(h or "").strip() for h in table.get("headers", [])]
        if len(headers) < 2:
            continue
        low_headers = [h.lower() for h in headers]
        label_idx = 0
        for i, h in enumerate(low_headers):
            if any(x in h for x in label_headers):
                label_idx = i
                break
        value_idx = None
        for term in preferred_value_headers:
            for i, h in enumerate(low_headers):
                if i != label_idx and term in h:
                    value_idx = i
                    break
            if value_idx is not None:
                break
        if value_idx is None:
            # fall back to first numeric column after the label
            for i in range(len(headers)):
                if i == label_idx:
                    continue
                for row in table.get("rows", [])[:6]:
                    if i < len(row) and parse_number(row[i]) is not None:
                        value_idx = i
                        break
                if value_idx is not None:
                    break
        if value_idx is None:
            continue
        for row in table.get("rows", []) or []:
            if len(row) <= max(label_idx, value_idx):
                continue
            label = clean_cell(row[label_idx])
            value = clean_cell(row[value_idx])
            if not label or parse_number(value) is None:
                continue
            # Avoid turning financial-bridge or footnote rows into KPI cards.
            low_label = label.lower().strip()
            if low_label.startswith(("higher ", "lower ", "first quarter", "second quarter", "third quarter", "fourth quarter", "less:", "add:", "total ")):
                continue
            # Keep executive-friendly labels.
            if len(label.split()) > 6:
                continue
            key = (label.lower(), value.lower())
            if key in seen:
                continue
            cards.append({"label": label, "value": value})
            seen.add(key)
            if len(cards) >= max_cards:
                return cards
    return cards

# ─── main entry point ────────────────────────────────────────────

def _empty_extraction(path, ext="", error=""):
    return {
        "text": "",
        "tables": [],
        "numbers": [],
        "metric_cards": [],
        "chart_candidates": [],
        "sections": [],
        "file_type": ext,
        "text_length": 0,
        "table_count": 0,
        "metric_count": 0,
        "chart_candidate_count": 0,
        "extraction_error": error,
        "source_file": path,
    }


def extract_document_data(path):
    ext = os.path.splitext(path)[1].lower()

    try:
        if ext == ".pdf":
            text, tables = extract_pdf(path)

        elif ext == ".docx":
            text, tables = extract_docx(path)

        elif ext == ".doc":
            try:
                import subprocess, tempfile, shutil
                tmp = tempfile.mkdtemp()
                subprocess.run(
                    [
                        "python3",
                        "/mnt/skills/public/docx/scripts/office/soffice.py",
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        tmp,
                        path,
                    ],
                    check=True,
                    capture_output=True,
                )
                converted = os.path.join(
                    tmp,
                    os.path.basename(path).replace(".doc", ".docx")
                )
                text, tables = extract_docx(converted)
                shutil.rmtree(tmp, ignore_errors=True)
            except Exception:
                with open(path, "rb") as f:
                    text = f.read().decode("utf-8", errors="ignore")
                tables = []

        elif ext == ".pptx":
            text, tables = extract_pptx(path)

        elif ext == ".txt":
            text, tables = extract_txt(path)

        elif ext == ".csv":
            text, tables = extract_csv(path)

        elif ext in (".xlsx", ".xls"):
            text, tables = extract_excel(path)

        else:
            return _empty_extraction(
                path,
                ext=ext,
                error=f"Unsupported file type: {ext}",
            )

        numbers = extract_numbers_from_text(text)
        metric_cards = deduplicate_metric_cards(
            extract_metric_cards(text) + extract_metric_cards_from_tables(tables)
        )

        chart_candidates = build_chart_candidates_from_tables(
            tables,
            metric_cards=metric_cards,
        )

        line_candidates = [] if _is_research_like_text(text) else build_chart_candidates_from_text_lines(text)

        seen_chart_titles = {
            str(c.get("title", "")).strip().lower()
            for c in chart_candidates
        }

        for candidate in line_candidates:
            key = str(candidate.get("title", "")).strip().lower()
            if key and key not in seen_chart_titles:
                chart_candidates.append(candidate)
                seen_chart_titles.add(key)

        sections = extract_sections(text)
        tables = normalize_single_column_tables(tables)

        return {
            "text": text,
            "tables": tables,
            "numbers": numbers,
            "metric_cards": metric_cards,
            "chart_candidates": chart_candidates,
            "sections": sections,
            "file_type": ext,
            "text_length": len(text),
            "table_count": len(tables),
            "metric_count": len(metric_cards),
            "chart_candidate_count": len(chart_candidates),
            "extraction_error": "",
            "source_file": path,
        }

    except Exception as exc:
        print(f"[data_extractor] extraction failed for {path}: {exc}")
        return _empty_extraction(path, ext=ext, error=str(exc))

def deduplicate_metric_cards(cards):
    """
    Keep only the best (most specific) card per label.
    For revenue/bookings that appear multiple times, keep the largest value.
    """
    from collections import defaultdict
    groups = defaultdict(list)
    for card in cards:
        groups[card["label"].lower()].append(card)

    deduped = []
    for label_key, group in groups.items():
        if len(group) == 1:
            deduped.append(group[0])
            continue
        # For financial metrics, prefer the largest value
        best = group[0]
        best_val = parse_number(best["value"]) or 0
        for g in group[1:]:
            gv = parse_number(g["value"]) or 0
            if gv > best_val:
                best = g
                best_val = gv
        deduped.append(best)
    return deduped
